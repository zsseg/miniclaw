"""高级文档生成工具：支持生成 PDF / PPT / 图文合成文档。

功能：
1. pdf  —— 用 fpdf2 将 Markdown/纯文本转为带格式的 PDF
2. ppt  —— 用 python-pptx 生成 PPT 幻灯片
3. compose —— 合图+文稿：先调用 image_gen 生成图片，再调用 writing_tool 生成文稿，最后打包为 PDF
"""

from __future__ import annotations

import json
import os
import re
import subprocess
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any, Callable

from fpdf import FPDF  # fpdf2
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from clawmini.tools.base import BaseTool
from clawmini.tools.registry import tool_plugin
from clawmini.types import ToolResult

# ── Unicode 字体探测 ────────────────────────
_CJK_FONT_CANDIDATES = [
    r"C:\Windows\Fonts\simhei.ttf",   # 黑体 — 纯 TTF，fpdf2 兼容最好
    r"C:\Windows\Fonts\Deng.ttf",     # 等线
    r"C:\Windows\Fonts\msyh.ttc",     # 微软雅黑 (TTC，部分 fpdf2 版本受限)
    r"C:\Windows\Fonts\simsun.ttc",   # 宋体
    "/System/Library/Fonts/PingFang.ttc",  # macOS
    "/System/Library/Fonts/STHeiti Light.ttc",
    "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",  # Linux
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
]


def _find_cjk_font() -> str:
    # 优先返回 .ttf 文件（fpdf2 对 ttf 支持比 ttc 好）
    ttf_candidates = [fp for fp in _CJK_FONT_CANDIDATES if fp.lower().endswith(".ttf")]
    ttc_candidates = [fp for fp in _CJK_FONT_CANDIDATES if fp.lower().endswith(".ttc")]
    for fp in ttf_candidates + ttc_candidates:
        if os.path.exists(fp):
            return fp
    # fallback: 在 Windows Fonts 目录下找任意 .ttf
    for root_dir in [r"C:\Windows\Fonts", "/usr/share/fonts"]:
        if os.path.isdir(root_dir):
            for fname in sorted(os.listdir(root_dir)):
                if fname.lower().endswith(".ttf"):
                    full = os.path.join(root_dir, fname)
                    if os.path.exists(full):
                        return full
            for fname in sorted(os.listdir(root_dir)):
                if fname.lower().endswith(".ttc"):
                    full = os.path.join(root_dir, fname)
                    if os.path.exists(full):
                        return full
    return ""


_CJK_FONT_PATH = _find_cjk_font()
_USE_UNICODE = bool(_CJK_FONT_PATH)


def _pdf_font(pdf: FPDF, style: str = "", size: int = 12) -> None:
    """设置 PDF 字体：有中文字体时用 Unicode 字体，否则用 Helvetica。

    注意：fpdf2 对 CJK 字体只注册了常规字重（style=''），
    因此即使传入 style='B' 也使用常规字体。
    """
    effective_style = "" if _USE_UNICODE else style
    if _USE_UNICODE:
        pdf.set_font("CJK", effective_style, size)
    else:
        pdf.set_font("Helvetica", style, size)


# ──────────────────────────────────────────────
# PDF 生成（基于 fpdf2）
# ──────────────────────────────────────────────

class _ClawminiPDF(FPDF):
    """自定义 PDF 子类：页眉/页脚，支持中文。"""

    def __init__(self, title: str = "") -> None:
        super().__init__()
        self._doc_title = title
        self.set_auto_page_break(auto=True, margin=20)
        if _USE_UNICODE:
            self.add_font("CJK", "", _CJK_FONT_PATH, uni=True)

    def header(self) -> None:
        if self._doc_title:
            _pdf_font(self, "", 8)
            self.set_text_color(140, 140, 140)
            self.cell(0, 8, self._doc_title, align="C", new_x="LMARGIN", new_y="NEXT")
            self.line(10, self.get_y(), 200, self.get_y())
            self.ln(4)

    def footer(self) -> None:
        self.set_y(-15)
        _pdf_font(self, "", 8)
        self.set_text_color(160, 160, 160)
        self.cell(0, 10, f"第 {self.page_no()}/{{nb}} 页", align="C")


def _markdown_simple_to_pdf(text: str, title: str = "", font_size: int = 12) -> _ClawminiPDF:
    """将 Markdown 简化文本写入 PDF 对象（支持 # ## ** - 等基础语法）。"""
    pdf = _ClawminiPDF(title=title)
    pdf.alias_nb_pages()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=20)

    lines = text.split("\n")
    for raw_line in lines:
        line = raw_line.rstrip()
        if not line:
            pdf.ln(4)
            continue

        # 标题 #
        h_match = re.match(r"^(#{1,3})\s+(.+)$", line)
        if h_match:
            level = len(h_match.group(1))
            content = h_match.group(2)
            sizes = {1: 22, 2: 18, 3: 14}
            _pdf_font(pdf, "B", sizes.get(level, 14))
            if level == 1:
                pdf.set_text_color(26, 60, 92)
            elif level == 2:
                pdf.set_text_color(50, 90, 130)
            else:
                pdf.set_text_color(70, 100, 140)
            pdf.cell(0, 10, content, new_x="LMARGIN", new_y="NEXT")
            if level == 1:
                pdf.line(10, pdf.get_y(), 200, pdf.get_y())
                pdf.ln(4)
            continue

        # 列表项 -
        list_match = re.match(r"^[-*]\s+(.+)$", line)
        if list_match:
            _pdf_font(pdf, "", font_size)
            pdf.set_text_color(40, 40, 40)
            pdf.set_x(pdf.l_margin + 4)
            pdf.multi_cell(pdf.w - pdf.l_margin - pdf.r_margin - 4, 6, "- " + list_match.group(1))
            continue

        # 数字列表
        num_match = re.match(r"^\d+[.．、]\s+(.+)$", line)
        if num_match:
            _pdf_font(pdf, "", font_size)
            pdf.set_text_color(40, 40, 40)
            pdf.set_x(pdf.l_margin + 4)
            pdf.multi_cell(pdf.w - pdf.l_margin - pdf.r_margin - 4, 6, "- " + num_match.group(1))
            continue

        # 普通段落
        _pdf_font(pdf, "", font_size)
        pdf.set_text_color(40, 40, 40)
        pdf.set_x(pdf.l_margin)
        pdf.multi_cell(0, 6, line)

    return pdf


# ──────────────────────────────────────────────
# PPT 生成（基于 python-pptx）
# ──────────────────────────────────────────────

def _text_to_ppt(
    text: str,
    title: str = "",
) -> Presentation:
    """将 Markdown 简化文本转换为 PPT 幻灯片。

    - # 一级标题 → 封面页
    - ## 二级标题 → 新幻灯片（标题+正文）
    - ### 三级标题 → 幻灯片内加粗小节
    - 其余内容作为正文段落
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    lines = text.split("\n")
    current_slide = None
    current_title = ""
    body_lines: list[str] = []
    first = True

    def _flush_slide() -> None:
        nonlocal body_lines, current_title, current_slide
        if current_slide is None and body_lines:
            # 封面用空白布局
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(26, 60, 92)
            # 标题
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = current_title or "文档标题"
            p.font.size = Pt(44)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            # 正文区
            if body_lines:
                txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9), Inches(2))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                for i, bline in enumerate(body_lines):
                    if i == 0:
                        p2 = tf2.paragraphs[0]
                    else:
                        p2 = tf2.add_paragraph()
                    p2.text = bline
                    p2.font.size = Pt(18)
                    p2.font.color.rgb = RGBColor(220, 230, 240)
                    p2.alignment = PP_ALIGN.CENTER
        elif current_slide is not None:
            # 内容页
            content = "\n".join(body_lines) if body_lines else ""
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # title and content
            if slide.shapes.title:
                slide.shapes.title.text = current_title
            if content and slide.placeholders:
                try:
                    ph = slide.placeholders[1]
                    tf = ph.text_frame
                    tf.word_wrap = True
                    tf.text = content
                except Exception:
                    pass

        body_lines = []

    for raw_line in lines:
        line = raw_line.rstrip()
        h1_match = re.match(r"^#\s+(.+)$", line)
        h2_match = re.match(r"^##\s+(.+)$", line)
        h3_match = re.match(r"^###\s+(.+)$", line)

        if h1_match:
            # 刷新上一张
            if current_slide is not None or body_lines:
                _flush_slide()
            current_slide = None
            current_title = h1_match.group(1)
            body_lines = []
            first = False
        elif h2_match:
            if current_slide is not None or body_lines:
                _flush_slide()
            current_slide = True
            current_title = h2_match.group(1)
            body_lines = []
        elif h3_match:
            body_lines.append(f"【{h3_match.group(1)}】")
        elif line.strip():
            body_lines.append(line)
        else:
            if body_lines and body_lines[-1] != "":
                body_lines.append("")

    _flush_slide()
    return prs


# ──────────────────────────────────────────────
# 图文合成：图片+文稿→PDF
# ──────────────────────────────────────────────

def _compose_image_and_text_to_pdf(
    image_paths: list[str],
    text: str,
    title: str = "",
) -> _ClawminiPDF:
    """将图片和文稿合并为一本 PDF。

    顺序：封面(标题) → 文稿页 → 图片页(每图一页)
    """
    pdf = _ClawminiPDF(title=title)
    pdf.alias_nb_pages()

    # 封面
    pdf.add_page()
    _pdf_font(pdf, "B", 28)
    pdf.set_text_color(26, 60, 92)
    pdf.ln(50)
    pdf.cell(0, 20, title or "图文报告", align="C", new_x="LMARGIN", new_y="NEXT")
    _pdf_font(pdf, "", 12)
    pdf.set_text_color(100, 100, 100)
    pdf.ln(10)
    pdf.cell(0, 10, f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(10)
    pdf.cell(0, 10, f"图片数：{len(image_paths)}", align="C", new_x="LMARGIN", new_y="NEXT")

    # 文稿内容
    if text.strip():
        pdf.add_page()
        _pdf_font(pdf, "B", 18)
        pdf.set_text_color(26, 60, 92)
        pdf.cell(0, 12, "文稿内容", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(4)
        # 写入文稿 — 注意 fpdf2 v2.5+ 的 multi_cell 不会自动重置 x 到左边距
        _pdf_font(pdf, "", 11)
        pdf.set_text_color(40, 40, 40)
        for para_line in text.split("\n"):
            pdf.set_x(pdf.l_margin)
            if para_line.strip():
                pdf.multi_cell(0, 6, para_line)
            else:
                pdf.ln(3)

    # 图片页
    for img_path_str in image_paths:
        img_path = Path(img_path_str)
        if not img_path.exists():
            continue
        try:
            pdf.add_page()
            # 图片标题
            _pdf_font(pdf, "B", 12)
            pdf.set_text_color(26, 60, 92)
            pdf.cell(0, 10, f"图片：{img_path.name}", new_x="LMARGIN", new_y="NEXT")
            # 尝试插入图片（适配页面大小）
            page_w = pdf.w - 2 * pdf.l_margin
            pdf.image(str(img_path), x=pdf.l_margin, y=pdf.get_y() + 4, w=page_w)
        except Exception:
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(180, 60, 60)
            pdf.cell(0, 10, f"[图片插入失败：{img_path.name}]", new_x="LMARGIN", new_y="NEXT")

    return pdf


# ──────────────────────────────────────────────
# 工具注册
# ──────────────────────────────────────────────

@tool_plugin("document_generation")
class DocumentGenerationTool(BaseTool):
    """高级文档生成工具：生成 PDF、PPT、图文合成文档。"""

    name = "document_generation"
    description = "生成 PDF/PPT/图文报告：指定 format='pdf'/'ppt'/'compose'，传入 content/title/image_paths 等"
    parameters_schema: dict[str, Any] = {
        "type": "object",
        "properties": {
            "command": {
                "type": "string",
                "description": "操作类型：create(生成)/info(查询信息)",
            },
            "format": {
                "type": "string",
                "description": "文档格式：pdf / ppt / compose",
            },
            "content": {
                "type": "string",
                "description": "文档内容（支持 Markdown 基础语法：标题/列表/加粗）",
            },
            "title": {
                "type": "string",
                "description": "文档标题",
            },
            "filename": {
                "type": "string",
                "description": "输出文件名（不含路径，后缀自动补全）",
            },
            "output_dir": {
                "type": "string",
                "description": "输出目录（工作区内相对路径）",
            },
            "image_paths": {
                "type": "array",
                "items": {"type": "string"},
                "description": "compose 模式需要的图片路径列表",
            },
            "font_size": {
                "type": "integer",
                "description": "PDF 正文字号（默认 12）",
            },
        },
        "required": ["command", "format"],
    }

    def __init__(self, workspace_dir: Path) -> None:
        super().__init__(workspace_dir)
        self._output_base = workspace_dir / "generated_documents"
        self._output_base.mkdir(parents=True, exist_ok=True)

    def run(self, arguments: dict[str, Any], progress_callback: Callable[[str], None] | None = None) -> ToolResult:
        command = str(arguments.get("command", "")).strip().lower()
        fmt = str(arguments.get("format", "pdf")).strip().lower()
        content = str(arguments.get("content", "")).strip()
        title = str(arguments.get("title", "")).strip() or "未命名文档"
        filename = str(arguments.get("filename", "")).strip()
        output_dir_str = str(arguments.get("output_dir", "")).strip()
        image_paths = arguments.get("image_paths", [])
        font_size = int(arguments.get("font_size", 12))

        if command == "info":
            return ToolResult(
                success=True,
                output=(
                    f"📄 文档生成工具\n"
                    f"支持的格式：pdf / ppt / compose(图片+文稿→PDF)\n"
                    f"输出目录：{self._output_base}\n"
                    f"pdf: 基于 fpdf2，支持 Markdown 基础语法\n"
                    f"ppt: 基于 python-pptx，支持 16:9 宽屏\n"
                    f"compose: 图片+文稿合并为 PDF"
                ),
            )

        if command != "create":
            return ToolResult(success=False, output=f"不支持的 command：{command}")

        if not content and fmt != "compose":
            return ToolResult(success=False, output="请提供 content（文档内容）")

        # 解析输出路径
        output_dir = self._output_base
        if output_dir_str:
            candidate = Path(output_dir_str)
            if not candidate.is_absolute():
                candidate = self.workspace_dir / candidate
            output_dir = candidate
        output_dir.mkdir(parents=True, exist_ok=True)

        # 自动确定文件名
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        if not filename:
            safe_title = re.sub(r"[\\/:*?\"<>|]+", "_", title)[:40]
            filename = f"{safe_title}_{ts}"
        else:
            filename = re.sub(r"[\\/:*?\"<>|]+", "_", filename)

        try:
            if fmt == "pdf":
                return self._generate_pdf(content, title, filename, output_dir, font_size)
            elif fmt == "ppt":
                return self._generate_ppt(content, title, filename, output_dir)
            elif fmt == "compose":
                return self._generate_compose(image_paths, content, title, filename, output_dir)
            else:
                return ToolResult(success=False, output=f"不支持的格式：{fmt}，可选：pdf / ppt / compose")
        except Exception as exc:
            return ToolResult(success=False, output=f"文档生成失败：{exc}")

    def _generate_pdf(
        self, content: str, title: str, filename: str, output_dir: Path, font_size: int
    ) -> ToolResult:
        pdf = _markdown_simple_to_pdf(content, title=title, font_size=font_size)
        pdf_path = output_dir / f"{filename}.pdf"
        pdf.output(str(pdf_path))
        pages = pdf.pages_count
        word_count = len(content.replace("\n", ""))
        return ToolResult(
            success=True,
            output=(
                f"✅ PDF 已生成！\n"
                f"文件：{pdf_path}\n"
                f"页数：{pages}\n"
                f"字数：{word_count}\n"
                f"标题：{title}"
            ),
            meta={
                "format": "pdf",
                "path": str(pdf_path),
                "pages": pages,
                "word_count": word_count,
                "title": title,
            },
        )

    def _generate_ppt(
        self, content: str, title: str, filename: str, output_dir: Path
    ) -> ToolResult:
        prs = _text_to_ppt(content, title=title)
        ppt_path = output_dir / f"{filename}.pptx"
        prs.save(str(ppt_path))
        slide_count = len(prs.slides)
        word_count = len(content.replace("\n", ""))
        return ToolResult(
            success=True,
            output=(
                f"✅ PPT 已生成！\n"
                f"文件：{ppt_path}\n"
                f"幻灯片数：{slide_count}\n"
                f"字数：{word_count}\n"
                f"标题：{title}"
            ),
            meta={
                "format": "ppt",
                "path": str(ppt_path),
                "slides": slide_count,
                "word_count": word_count,
                "title": title,
            },
        )

    def _generate_compose(
        self, image_paths: list[str], content: str, title: str, filename: str, output_dir: Path
    ) -> ToolResult:
        # 解析图片路径
        resolved_images: list[str] = []
        for p in image_paths:
            pp = Path(p)
            if not pp.is_absolute():
                pp = self.workspace_dir / pp
            if pp.exists():
                resolved_images.append(str(pp.resolve()))

        pdf = _compose_image_and_text_to_pdf(resolved_images, content, title=title)
        pdf_path = output_dir / f"{filename}.pdf"
        pdf.output(str(pdf_path))
        pages = pdf.pages_count
        word_count = len(content.replace("\n", ""))
        return ToolResult(
            success=True,
            output=(
                f"✅ 图文合成 PDF 已生成！\n"
                f"文件：{pdf_path}\n"
                f"页数：{pages}\n"
                f"图片数：{len(resolved_images)}\n"
                f"字数：{word_count}\n"
                f"标题：{title}\n"
                f"图片列表：\n" + "\n".join(resolved_images)
            ),
            meta={
                "format": "compose",
                "path": str(pdf_path),
                "pages": pages,
                "images": resolved_images,
                "word_count": word_count,
                "title": title,
            },
        )
